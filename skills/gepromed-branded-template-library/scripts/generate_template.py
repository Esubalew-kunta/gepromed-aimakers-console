#!/usr/bin/env python3
"""Generate a GEPROMED-branded TEMPLATE: letter (.docx), report (.docx),
presentation (.pptx), or email (.txt).

Deterministic renderer applying the GEPROMED charte (blue #007AC2, rare orange
#EC6C17, dark #1F2A33, muted #5F6B73, bundled logo, Calibri). It lays out the
structure and fills it with the author's supplied content. It does NOT invent
facts, figures, or proof points — unknown values stay as `[crochets]` for the
author to confirm. Follows references/template-specs.md.

Input is a JSON file describing the document (see --print-schema) or use --demo.

Usage:
    python generate_template.py --type letter       --in doc.json --out letter.docx
    python generate_template.py --type report       --in doc.json --out report.docx
    python generate_template.py --type presentation --in doc.json --out deck.pptx
    python generate_template.py --type email        --in doc.json --out email.txt
    python generate_template.py --type report --demo --out demo_report.docx
    python generate_template.py --print-schema

Exit 0 on success, 2 on bad input.
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from docx import Document
from docx.shared import Pt, RGBColor, Inches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

SCRIPT_DIR = Path(__file__).resolve().parent
LOGO = SCRIPT_DIR.parent / "assets" / "gepromed-logo.png"

BLUE = RGBColor(0x00, 0x7A, 0xC2)
ORANGE = RGBColor(0xEC, 0x6C, 0x17)
DARK = RGBColor(0x1F, 0x2A, 0x33)
MUTED = RGBColor(0x5F, 0x6B, 0x73)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Calibri"

# pptx colours
from pptx.util import Pt as PPt, Inches as PInches, Emu  # noqa: E402
from pptx.dml.color import RGBColor as PRGB  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402

P_BLUE = PRGB(0x00, 0x7A, 0xC2)
P_ORANGE = PRGB(0xEC, 0x6C, 0x17)
P_DARK = PRGB(0x1F, 0x2A, 0x33)
P_MUTED = PRGB(0x5F, 0x6B, 0x73)
P_WHITE = PRGB(0xFF, 0xFF, 0xFF)


def _ph(value, placeholder: str) -> str:
    return str(value) if value not in (None, "", []) else placeholder


# ----------------------------- docx helpers -----------------------------

def _base_styles(doc: Document) -> None:
    normal = doc.styles["Normal"]
    normal.font.name = FONT
    normal.font.size = Pt(10.5)
    normal.font.color.rgb = DARK


def _logo(doc: Document, width=2.0) -> None:
    if LOGO.exists():
        p = doc.add_paragraph()
        p.add_run().add_picture(str(LOGO), width=Inches(width))
    t = doc.add_paragraph()
    tr = t.add_run("The medical device hub for patient safety")
    tr.italic = True
    tr.font.size = Pt(8.5)
    tr.font.color.rgb = MUTED
    t.paragraph_format.space_after = Pt(6)


def _heading(doc: Document, text: str, size=13) -> None:
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.bold = True
    run.font.size = Pt(size)
    run.font.color.rgb = BLUE
    pPr = p._p.get_or_add_pPr()
    pbdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), "6")
    bottom.set(qn("w:space"), "2")
    bottom.set(qn("w:color"), "007AC2")
    pbdr.append(bottom)
    pPr.append(pbdr)


def _body(doc: Document, text: str, muted=False) -> None:
    p = doc.add_paragraph()
    r = p.add_run(text)
    r.font.size = Pt(10.5)
    r.font.color.rgb = MUTED if muted else DARK
    p.paragraph_format.space_after = Pt(4)


def _bullets(doc: Document, items) -> None:
    for it in items:
        p = doc.add_paragraph(style="List Bullet")
        r = p.add_run(str(it))
        r.font.size = Pt(10.5)
        r.font.color.rgb = DARK


def _doc_footer(doc: Document, text: str) -> None:
    p = doc.sections[0].footer.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(text)
    r.font.size = Pt(7.5)
    r.font.color.rgb = MUTED


def _render_sections(doc: Document, sections) -> None:
    """sections: list of {title, body?(str), bullets?(list)} or plain strings."""
    for s in sections:
        if isinstance(s, str):
            _body(doc, s)
            continue
        if s.get("title"):
            _heading(doc, s["title"])
        if s.get("body"):
            _body(doc, s["body"])
        if s.get("bullets"):
            _bullets(doc, s["bullets"])


# ----------------------------- letter -----------------------------

SALUTATIONS = {
    "fr": {
        "hcp": ("Cher Docteur,", "Bien cordialement,"),
        "industry": ("Bonjour,", "Cordialement,"),
        "researcher": ("Cher collègue,", "Bien cordialement,"),
        "institution": ("Madame, Monsieur,", "Je vous prie d'agréer, Madame, Monsieur, l'expression de mes salutations distinguées."),
        "participant": ("Bonjour,", "Bien cordialement,"),
        "internal": ("Bonjour,", "Merci,"),
        "default": ("Madame, Monsieur,", "Bien cordialement,"),
    },
    "en": {
        "hcp": ("Dear Dr [Name],", "Kind regards,"),
        "industry": ("Dear [First name],", "Best regards,"),
        "researcher": ("Dear [First name],", "With best regards,"),
        "institution": ("Dear Sir or Madam,", "Yours faithfully,"),
        "participant": ("Hello [First name],", "Kind regards,"),
        "internal": ("Hi [First name],", "Thanks,"),
        "default": ("Dear Sir or Madam,", "Best regards,"),
    },
}


def build_letter(data: dict, out: Path) -> Path:
    fr = (data.get("language", "fr").lower() != "en")
    lang = "fr" if fr else "en"
    aud = data.get("audience", "default")
    sal, clo = SALUTATIONS[lang].get(aud, SALUTATIONS[lang]["default"])

    doc = Document()
    for s in doc.sections:
        s.top_margin = Cm(2.0); s.bottom_margin = Cm(2.0)
        s.left_margin = Cm(2.5); s.right_margin = Cm(2.5)
    _base_styles(doc)
    _logo(doc)

    # recipient block
    rec = data.get("recipient", {})
    rp = doc.add_paragraph()
    rp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    rec_lines = [
        _ph(rec.get("name"), "[Nom du destinataire]" if fr else "[Recipient name]"),
        _ph(rec.get("organisation"), "[Organisation]"),
        _ph(rec.get("address"), "[Adresse]" if fr else "[Address]"),
    ]
    rr = rp.add_run("\n".join(rec_lines))
    rr.font.size = Pt(10)
    rr.font.color.rgb = DARK

    # place + date
    dp = doc.add_paragraph()
    dp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    dr = dp.add_run(("Strasbourg, le " if fr else "Strasbourg, ") + _ph(data.get("date"), "[date]"))
    dr.font.size = Pt(10)
    dr.font.color.rgb = MUTED
    dp.paragraph_format.space_after = Pt(10)

    # objet
    op = doc.add_paragraph()
    olab = "Objet : " if fr else "Subject: "
    olr = op.add_run(olab)
    olr.bold = True; olr.font.color.rgb = BLUE; olr.font.size = Pt(11)
    ovr = op.add_run(_ph(data.get("subject") or data.get("purpose"),
                         "[objet du courrier]" if fr else "[letter subject]"))
    ovr.bold = True; ovr.font.color.rgb = BLUE; ovr.font.size = Pt(11)
    op.paragraph_format.space_after = Pt(8)

    # salutation
    _body(doc, sal)
    # body sections
    _render_sections(doc, data.get("sections", []))
    if not data.get("sections"):
        _body(doc, "[Corps du courrier — paragraphes à compléter.]" if fr else
                   "[Letter body — paragraphs to complete.]", muted=True)
    # closing
    doc.add_paragraph()
    _body(doc, clo)
    sig = doc.add_paragraph()
    sname = _ph(data.get("author_name"), "[Prénom Nom]" if fr else "[First name Last name]")
    srole = _ph(data.get("author_role"), "[Fonction]" if fr else "[Role]")
    sr = sig.add_run(f"{sname}\n{srole}\nGEPROMED — The medical device hub for patient safety")
    sr.font.size = Pt(10)
    sr.font.color.rgb = DARK

    _doc_footer(doc, ("GEPROMED · Strasbourg · gepromed.com   |   Document de travail — relire avant envoi.") if fr else
                     ("GEPROMED · Strasbourg · gepromed.com   |   Working document — review before sending."))
    out.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out))
    return out


# ----------------------------- report -----------------------------

def build_report(data: dict, out: Path) -> Path:
    fr = (data.get("language", "fr").lower() != "en")
    doc = Document()
    for s in doc.sections:
        s.top_margin = Cm(2.0); s.bottom_margin = Cm(2.0)
        s.left_margin = Cm(2.2); s.right_margin = Cm(2.2)
    _base_styles(doc)
    _logo(doc)

    # title
    tp = doc.add_paragraph()
    tr = tp.add_run(_ph(data.get("title"), "[Titre du rapport]" if fr else "[Report title]"))
    tr.bold = True; tr.font.size = Pt(20); tr.font.color.rgb = DARK
    tag = doc.add_paragraph()
    tgr = tag.add_run("Rapport" if fr else "Report")
    tgr.bold = True; tgr.font.size = Pt(11); tgr.font.color.rgb = ORANGE
    # meta
    meta_bits = []
    if data.get("reference"):
        meta_bits.append((f"Réf. {data['reference']}") if fr else (f"Ref. {data['reference']}"))
    meta_bits.append(f"Version {data.get('version','[v.]')}")
    meta_bits.append((f"Date : {data.get('date','[date]')}") if fr else (f"Date: {data.get('date','[date]')}"))
    meta_bits.append((f"Auteur : {data.get('author','[auteur]')}") if fr else (f"Author: {data.get('author','[author]')}"))
    mp = doc.add_paragraph()
    mr = mp.add_run("  ·  ".join(meta_bits))
    mr.font.size = Pt(8.5); mr.font.color.rgb = MUTED
    mp.paragraph_format.space_after = Pt(10)

    _render_sections(doc, data.get("sections", []))
    if not data.get("sections"):
        _heading(doc, "Synthèse" if fr else "Executive summary")
        _body(doc, "[À compléter.]" if fr else "[To complete.]", muted=True)

    _doc_footer(doc, ("GEPROMED · Strasbourg · gepromed.com   |   Document de travail — relire avant diffusion.") if fr else
                     ("GEPROMED · Strasbourg · gepromed.com   |   Working document — review before release."))
    out.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out))
    return out


# ----------------------------- presentation -----------------------------

def build_presentation(data: dict, out: Path) -> Path:
    from pptx import Presentation

    fr = (data.get("language", "fr").lower() != "en")
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]
    SW, SH = prs.slide_width, prs.slide_height

    def add_logo(slide, top=PInches(6.7)):
        if LOGO.exists():
            slide.shapes.add_picture(str(LOGO), PInches(0.4), top, width=PInches(1.6))

    def textbox(slide, left, top, width, height):
        tb = slide.shapes.add_textbox(left, top, width, height)
        return tb.text_frame

    # ---- title slide
    s = prs.slides.add_slide(blank)
    band = s.shapes.add_shape(1, 0, 0, SW, PInches(2.6))
    band.fill.solid(); band.fill.fore_color.rgb = P_BLUE
    band.line.fill.background()
    accent = s.shapes.add_shape(1, 0, PInches(2.6), SW, PInches(0.08))
    accent.fill.solid(); accent.fill.fore_color.rgb = P_ORANGE
    accent.line.fill.background()
    tf = textbox(s, PInches(0.7), PInches(0.7), PInches(11.9), PInches(1.8))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = _ph(data.get("title"), "[Titre de la présentation]" if fr else "[Presentation title]")
    r.font.size = PPt(36); r.font.bold = True; r.font.color.rgb = P_WHITE
    sp = tf.add_paragraph()
    sr = sp.add_run(); sr.text = _ph(data.get("subtitle"), "[sous-titre]" if fr else "[subtitle]")
    sr.font.size = PPt(18); sr.font.color.rgb = P_WHITE
    mtf = textbox(s, PInches(0.7), PInches(3.0), PInches(11.9), PInches(0.6))
    mp = mtf.paragraphs[0]
    mr = mp.add_run()
    mr.text = _ph(data.get("meta"), "GEPROMED — The medical device hub for patient safety")
    mr.font.size = PPt(12); mr.font.color.rgb = P_MUTED
    add_logo(s)

    sections = data.get("sections", [])

    # ---- agenda slide
    if sections:
        a = prs.slides.add_slide(blank)
        atf = textbox(a, PInches(0.7), PInches(0.5), PInches(11.9), PInches(0.9))
        ap = atf.paragraphs[0]
        ar = ap.add_run(); ar.text = "Sommaire" if fr else "Agenda"
        ar.font.size = PPt(28); ar.font.bold = True; ar.font.color.rgb = P_BLUE
        btf = textbox(a, PInches(0.9), PInches(1.6), PInches(11.5), PInches(5.0))
        first = True
        for sec in sections:
            title = sec.get("title") if isinstance(sec, dict) else str(sec)
            para = btf.paragraphs[0] if first else btf.add_paragraph()
            first = False
            run = para.add_run(); run.text = "•  " + (title or "[section]")
            run.font.size = PPt(18); run.font.color.rgb = P_DARK
        add_logo(a)

    # ---- content slides
    for sec in sections:
        c = prs.slides.add_slide(blank)
        title = sec.get("title") if isinstance(sec, dict) else str(sec)
        ttf = textbox(c, PInches(0.7), PInches(0.5), PInches(11.9), PInches(0.9))
        tp = ttf.paragraphs[0]
        tr = tp.add_run(); tr.text = title or "[Titre de section]"
        tr.font.size = PPt(26); tr.font.bold = True; tr.font.color.rgb = P_BLUE
        rule = c.shapes.add_shape(1, PInches(0.7), PInches(1.35), PInches(3.0), PInches(0.05))
        rule.fill.solid(); rule.fill.fore_color.rgb = P_ORANGE; rule.line.fill.background()

        body_tf = textbox(c, PInches(0.9), PInches(1.7), PInches(11.5), PInches(4.8))
        bullets = []
        if isinstance(sec, dict):
            if sec.get("body"):
                bullets.append(sec["body"])
            bullets.extend(sec.get("bullets", []) or [])
        if not bullets:
            bullets = ["[Contenu à compléter.]" if fr else "[Content to complete.]"]
        first = True
        for b in bullets[:6]:
            para = body_tf.paragraphs[0] if first else body_tf.add_paragraph()
            first = False
            run = para.add_run(); run.text = "•  " + str(b)
            run.font.size = PPt(18); run.font.color.rgb = P_DARK
            para.space_after = PPt(8)
        add_logo(c)

    # ---- closing slide
    z = prs.slides.add_slide(blank)
    band = z.shapes.add_shape(1, 0, PInches(2.8), SW, PInches(1.9))
    band.fill.solid(); band.fill.fore_color.rgb = P_BLUE; band.line.fill.background()
    ztf = textbox(z, PInches(0.7), PInches(3.1), PInches(11.9), PInches(1.4))
    zp = ztf.paragraphs[0]
    zr = zp.add_run(); zr.text = "Merci" if fr else "Thank you"
    zr.font.size = PPt(40); zr.font.bold = True; zr.font.color.rgb = P_WHITE
    ctf = textbox(z, PInches(0.7), PInches(4.9), PInches(11.9), PInches(0.8))
    cp = ctf.paragraphs[0]
    cr = cp.add_run()
    cr.text = _ph(data.get("contact"), "[contact] · gepromed.com")
    cr.font.size = PPt(14); cr.font.color.rgb = P_MUTED
    add_logo(z)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


# ----------------------------- email -----------------------------

def build_email(data: dict, out: Path) -> Path:
    fr = (data.get("language", "fr").lower() != "en")
    lang = "fr" if fr else "en"
    aud = data.get("audience", "default")
    sal, clo = SALUTATIONS[lang].get(aud, SALUTATIONS[lang]["default"])
    subject = _ph(data.get("subject") or data.get("purpose"),
                  "[objet de l'email]" if fr else "[email subject]")

    lines = []
    lines.append(("Objet : " if fr else "Subject: ") + subject)
    lines.append("")
    lines.append(sal)
    lines.append("")
    purpose = _ph(data.get("purpose"),
                  "[Phrase d'objet : pourquoi vous écrivez.]" if fr else
                  "[Purpose line: why you are writing.]")
    lines.append(purpose)
    lines.append("")
    sections = data.get("sections", [])
    if sections:
        for s in sections:
            if isinstance(s, dict):
                if s.get("title"):
                    lines.append(s["title"])
                if s.get("body"):
                    lines.append(s["body"])
                for b in s.get("bullets", []) or []:
                    lines.append(f"- {b}")
            else:
                lines.append(f"- {s}")
        lines.append("")
    else:
        lines.append("- [point 1]")
        lines.append("- [point 2]")
        lines.append("")
    ask = _ph(data.get("ask"),
              "[Demande explicite : l'action attendue + l'échéance.]" if fr else
              "[Explicit ask: the action you need + the deadline.]")
    lines.append(ask)
    lines.append("")
    lines.append(clo)
    lines.append("")
    lines.append(_ph(data.get("author_name"), "[Prénom Nom]" if fr else "[First name Last name]"))
    lines.append(_ph(data.get("author_role"), "[Fonction]" if fr else "[Role]"))
    lines.append("GEPROMED — The medical device hub for patient safety")
    lines.append("[téléphone]  ·  [email]  ·  gepromed.com" if fr else "[phone]  ·  [email]  ·  gepromed.com")

    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return out


# ----------------------------- driver -----------------------------

BUILDERS = {
    "letter": build_letter,
    "report": build_report,
    "presentation": build_presentation,
    "email": build_email,
}

SCHEMA = {
    "language": "fr | en (default fr)",
    "audience": "hcp | industry | researcher | institution | participant | internal | default",
    "title": "(report/presentation) document title",
    "subtitle": "(presentation) subtitle",
    "meta": "(presentation) meta line",
    "subject": "(letter/email) subject/objet",
    "purpose": "(email/letter) one-line purpose",
    "ask": "(email) explicit ask",
    "recipient": {"name": "", "organisation": "", "address": ""},
    "date": "date (optional)",
    "reference": "(report) ref (optional)",
    "version": "(report) version (optional)",
    "author": "(report) author (optional)",
    "author_name": "(letter/email) signer name (optional)",
    "author_role": "(letter/email) signer role (optional)",
    "contact": "(presentation) closing contact (optional)",
    "sections": "list of {title, body?, bullets?[]} or strings — the outline",
}

DEMOS = {
    "report": {
        "language": "fr",
        "title": "Bilan d'activité — Plateforme de test 2026",
        "reference": "GEP-RAP-2026-01",
        "version": "1.0",
        "date": "[date]",
        "author": "[auteur]",
        "sections": [
            {"title": "Synthèse", "body": "Ce rapport présente l'activité de la plateforme de test sur la période [période]. Les chiffres ci-dessous sont à confirmer par le responsable de plateforme."},
            {"title": "Activité de test", "bullets": ["Essais mécaniques réalisés : [N]", "Dispositifs cardiovasculaires évalués : [N]", "Référentiels appliqués : ISO 9001, ISO 13485"]},
            {"title": "Faits marquants", "body": "[À compléter avec les éléments validés.]"},
            {"title": "Perspectives", "bullets": ["[Axe 1]", "[Axe 2]"]},
        ],
    },
    "letter": {
        "language": "fr",
        "audience": "institution",
        "subject": "Demande de soutien — programme de formation par simulation",
        "recipient": {"name": "[Nom]", "organisation": "[Institution]", "address": "[Adresse]"},
        "date": "[date]",
        "author_name": "[Prénom Nom]",
        "author_role": "[Fonction]",
        "sections": [
            {"body": "Dans le cadre de sa mission d'amélioration de la sécurité des patients, GEPROMED développe un programme de formation chirurgicale par simulation au René Kieny Education Center."},
            {"body": "Nous sollicitons votre soutien pour étendre ce programme à [périmètre]. Un dossier détaillé est disponible sur demande."},
        ],
    },
    "presentation": {
        "language": "fr",
        "title": "GEPROMED — Plateforme de test",
        "subtitle": "Essais mécaniques et biomatériaux des dispositifs cardiovasculaires",
        "meta": "GEPROMED — The medical device hub for patient safety",
        "contact": "[contact] · gepromed.com",
        "sections": [
            {"title": "Mission", "bullets": ["Sécurité des patients via le cycle de vie de l'implant", "Indépendance et rigueur scientifique"]},
            {"title": "Capacités de test", "bullets": ["Essais mécaniques", "Essais biomatériaux", "ISO 9001 · ISO 13485"]},
            {"title": "Collaboration", "bullets": ["[modalités à préciser]", "[contact]"]},
        ],
    },
    "email": {
        "language": "fr",
        "audience": "industry",
        "subject": "Analyse d'explants indépendante — GEPROMED",
        "purpose": "Je reviens vers vous au sujet d'une possible collaboration sur l'analyse d'explants.",
        "ask": "Seriez-vous disponible pour un court échange dans les prochaines semaines ?",
        "author_name": "[Prénom Nom]",
        "author_role": "[Fonction]",
        "sections": [{"bullets": ["Analyse d'explants indépendante", "Essais mécaniques et biomatériaux", "Cadre ISO 9001 / ISO 13485"]}],
    },
}


def main() -> int:
    ap = argparse.ArgumentParser(description="Generate a GEPROMED-branded template.")
    ap.add_argument("--type", choices=list(BUILDERS.keys()), help="Template type.")
    ap.add_argument("--in", dest="infile", help="Path to document JSON.")
    ap.add_argument("--out", help="Output path.")
    ap.add_argument("--print-schema", action="store_true")
    ap.add_argument("--demo", action="store_true")
    args = ap.parse_args()

    if args.print_schema:
        print(json.dumps(SCHEMA, ensure_ascii=False, indent=2))
        return 0

    if not args.type:
        print("ERROR: --type is required (letter | report | presentation | email).")
        return 2

    if args.demo:
        data = DEMOS[args.type]
    elif args.infile:
        try:
            data = json.loads(Path(args.infile).read_text(encoding="utf-8"))
        except Exception as e:  # noqa: BLE001
            print(f"ERROR reading {args.infile}: {e}")
            return 2
    else:
        print("ERROR: provide --in doc.json (or --demo / --print-schema).")
        return 2

    default_ext = {"letter": ".docx", "report": ".docx", "presentation": ".pptx", "email": ".txt"}
    out = Path(args.out) if args.out else Path(f"{args.type}{default_ext[args.type]}")
    result = BUILDERS[args.type](data, out)
    print(f"OK — wrote {result}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
