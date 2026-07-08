#!/usr/bin/env python3
"""Generate a brand-styled GEPROMED ISO 9001 §9.3 management-review deck (.pptx).

This is a DETERMINISTIC renderer. It does NOT invent content: it draws exactly what
the JSON content file (or the built-in --demo content) provides. Any value left as
a "[bracket ...]" string is rendered verbatim so the RQ can spot and fill it. The
script never generates KPI values, audit results, or decisions.

Styling: GEPROMED palette (blue master #007AC2, orange accent #EC6C17 used sparingly,
dark text #1F2A33, muted #5F6B73), the bundled logo on the title slide and every
footer, one slide per §9.3 section.

Requires: python-pptx  (pip install python-pptx)

Usage:
    python generate_review_deck.py --content content.json --out revue_direction.pptx
    python generate_review_deck.py --demo --out demo_revue_direction.pptx
    python generate_review_deck.py --print-schema

Exit codes: 0 ok · 2 bad input · 3 missing python-pptx.
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:  # pragma: no cover
    print(
        "ERROR: python-pptx is required. Install it with:\n    pip install python-pptx",
        file=sys.stderr,
    )
    raise SystemExit(3)

SCRIPT_DIR = Path(__file__).resolve().parent
DEFAULT_LOGO = SCRIPT_DIR.parent / "assets" / "gepromed-logo.png"

# --- GEPROMED palette ---
BLUE = RGBColor(0x00, 0x7A, 0xC2)
ORANGE = RGBColor(0xEC, 0x6C, 0x17)
DARK = RGBColor(0x1F, 0x2A, 0x33)
MUTED = RGBColor(0x5F, 0x6B, 0x73)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 16:9 canvas
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

SCHEMA = {
    "title": "Revue de direction — <période>",
    "period": "<période, ex. Année 2025>",
    "date": "<date de la réunion>",
    "audience": "<Direction / comité qualité / auditeur>",
    "language": "fr | en",
    "scope": "<périmètre du SMQ — [à confirmer par le RQ]>",
    "sections": [
        {
            "heading": "Section title (e.g. 'Satisfaction client')",
            "bullets": [
                "A supplied fact or a [bracketed — à confirmer par le RQ] placeholder",
                "Another bullet",
            ],
        }
    ],
    "decisions": ["Decision text or [à confirmer par le RQ]"],
    "actions": [
        {"action": "...", "owner": "[à attribuer]", "deadline": "[à définir]", "status": "[à faire]"}
    ],
    "validator_role": "RQ",
}

# Default §9.3 section order, all bracketed (used by --demo and as a fallback).
DEFAULT_SECTIONS = [
    ("Contexte & périmètre du SMQ", ["[périmètre du SMQ — à confirmer par le RQ]"]),
    ("Suivi des actions des revues précédentes", ["[statut par action — à confirmer par le RQ]"]),
    ("Évolutions internes & externes", ["[changements de contexte — à compléter par le RQ]"]),
    ("Satisfaction client & parties intéressées", ["[score de satisfaction — à confirmer]", "[retours qualitatifs — à confirmer]"]),
    ("Atteinte des objectifs qualité", ["[statut par objectif — à confirmer par le RQ]"]),
    ("Performance des processus", ["[KPI par processus — à confirmer par le RQ]"]),
    ("Non-conformités & actions correctives", ["[nombre & nature des NC — à confirmer]", "[statut des actions correctives — à confirmer]"]),
    ("Résultats d'audit", ["[constats / écarts d'audit — à confirmer par le RQ]"]),
    ("Prestataires externes", ["[performance des prestataires — à confirmer]"]),
    ("Ressources, risques & opportunités", ["[adéquation des ressources — à confirmer]", "[efficacité des actions risques/opportunités — à confirmer]"]),
    ("Opportunités d'amélioration", ["[opportunités identifiées — à compléter par le RQ]"]),
    ("Objectifs qualité — revus / nouveaux", ["[objectifs revus ou nouveaux — à confirmer par le RQ]"]),
]


def demo_content() -> dict:
    return {
        "title": "Revue de direction — Année [AAAA]",
        "period": "Année [AAAA] — [période à confirmer]",
        "date": "[date de la réunion]",
        "audience": "Direction & comité qualité",
        "language": "fr",
        "scope": "[périmètre du SMQ — à confirmer par le RQ]",
        "sections": [{"heading": h, "bullets": b} for h, b in DEFAULT_SECTIONS],
        "decisions": ["[décision 1 — à confirmer par le RQ]", "[décision 2 — à confirmer par le RQ]"],
        "actions": [
            {"action": "[action — à définir par le RQ]", "owner": "[responsable — à attribuer]",
             "deadline": "[échéance — à définir]", "status": "[à faire]"},
        ],
        "validator_role": "RQ",
    }


def _add_footer(slide, text: str, logo_path: Path | None) -> None:
    # Thin blue rule above the footer.
    line = slide.shapes.add_shape(1, Inches(0.0), Inches(7.02), SLIDE_W, Pt(2))
    line.fill.solid(); line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    # Footer text (muted).
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(9.5), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(9); r.font.color.rgb = MUTED
    # Small logo bottom-right.
    if logo_path and logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(11.7), Inches(7.0), height=Inches(0.35))


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _title_slide(prs, content, logo_path):
    slide = _blank(prs)
    # Blue band across the top.
    band = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(2.1))
    band.fill.solid(); band.fill.fore_color.rgb = BLUE; band.line.fill.background()
    # Orange accent strip (rare accent, thin).
    strip = slide.shapes.add_shape(1, 0, Inches(2.1), SLIDE_W, Pt(6))
    strip.fill.solid(); strip.fill.fore_color.rgb = ORANGE; strip.line.fill.background()
    # Logo centered upper area.
    if logo_path and logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(0.6), Inches(0.45), height=Inches(1.2))
    # Title.
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(12), Inches(1.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = content.get("title", "Revue de direction")
    r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = DARK
    # Subtitle line.
    sub = slide.shapes.add_textbox(Inches(0.7), Inches(3.9), Inches(12), Inches(2))
    tf = sub.text_frame
    meta = [
        ("Période / Period : ", content.get("period", "[période]")),
        ("Date : ", content.get("date", "[date]")),
        ("Audience : ", content.get("audience", "[audience]")),
        ("Référentiel : ", "ISO 9001 — §9.3 Revue de direction"),
        ("Statut : ", "PROJET — à valider par le RQ"),
    ]
    for i, (label, val) in enumerate(meta):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        rl = para.add_run(); rl.text = label
        rl.font.size = Pt(14); rl.font.bold = True; rl.font.color.rgb = MUTED
        rv = para.add_run(); rv.text = str(val)
        rv.font.size = Pt(14); rv.font.color.rgb = DARK
    _add_footer(slide, "GEPROMED — Revue de direction (projet) · à valider par le RQ", logo_path)
    return slide


def _section_header(slide, heading):
    # Blue title bar.
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = heading
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = WHITE
    # Small orange tab as a section marker (rare accent).
    tab = slide.shapes.add_shape(1, Inches(0.0), Inches(1.1), Inches(0.25), Inches(0.18))
    tab.fill.solid(); tab.fill.fore_color.rgb = ORANGE; tab.line.fill.background()


def _content_slide(prs, heading, bullets, logo_path):
    slide = _blank(prs)
    _section_header(slide, heading)
    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.2))
    tf = body.text_frame; tf.word_wrap = True
    if not bullets:
        bullets = ["[à compléter par le RQ]"]
    for i, b in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = para.add_run(); r.text = f"•  {b}"
        r.font.size = Pt(18)
        # Bracketed placeholders get the orange accent so the RQ spots them fast.
        r.font.color.rgb = ORANGE if str(b).strip().startswith("[") else DARK
        para.space_after = Pt(10)
    _add_footer(slide, "GEPROMED — Revue de direction (projet) · à valider par le RQ", logo_path)
    return slide


def _actions_slide(prs, decisions, actions, logo_path):
    slide = _blank(prs)
    _section_header(slide, "Décisions & plan d'actions")
    # Decisions block.
    db = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(1.6))
    tf = db.text_frame; tf.word_wrap = True
    head = tf.paragraphs[0]; hr = head.add_run(); hr.text = "Décisions :"
    hr.font.size = Pt(16); hr.font.bold = True; hr.font.color.rgb = BLUE
    for d in (decisions or ["[décision — à confirmer par le RQ]"]):
        p = tf.add_paragraph(); r = p.add_run(); r.text = f"•  {d}"
        r.font.size = Pt(13)
        r.font.color.rgb = ORANGE if str(d).strip().startswith("[") else DARK
    # Actions table.
    actions = actions or [{"action": "[à définir]", "owner": "[à attribuer]",
                           "deadline": "[à définir]", "status": "[à faire]"}]
    rows = len(actions) + 1
    cols = 5
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.7), Inches(3.3),
                                       Inches(12), Inches(0.4 * rows))
    table = tbl_shape.table
    headers = ["#", "Action", "Responsable", "Échéance", "Statut"]
    widths = [Inches(0.6), Inches(5.4), Inches(2.4), Inches(1.9), Inches(1.7)]
    for c, w in enumerate(widths):
        table.columns[c].width = w
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
        para = cell.text_frame.paragraphs[0]
        run = para.add_run(); run.text = h
        run.font.size = Pt(12); run.font.bold = True; run.font.color.rgb = WHITE
    for i, a in enumerate(actions, start=1):
        vals = [str(i), a.get("action", "[à définir]"), a.get("owner", "[à attribuer]"),
                a.get("deadline", "[à définir]"), a.get("status", "[à faire]")]
        for c, v in enumerate(vals):
            cell = table.cell(i, c)
            para = cell.text_frame.paragraphs[0]
            run = para.add_run(); run.text = v
            run.font.size = Pt(11)
            run.font.color.rgb = ORANGE if v.strip().startswith("[") else DARK
    _add_footer(slide, "GEPROMED — Revue de direction (projet) · à valider par le RQ", logo_path)
    return slide


def _closing_slide(prs, validator_role, logo_path):
    slide = _blank(prs)
    _section_header(slide, "Clôture & validation")
    box = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(3.5))
    tf = box.text_frame; tf.word_wrap = True
    p0 = tf.paragraphs[0]; r0 = p0.add_run()
    r0.text = "⚠️ VALIDATION — " + validator_role
    r0.font.size = Pt(22); r0.font.bold = True; r0.font.color.rgb = ORANGE
    p1 = tf.add_paragraph(); r1 = p1.add_run()
    r1.text = ("Ce document est un PROJET, pas un enregistrement qualité validé. "
               "Chaque KPI, résultat d'audit, non-conformité, score de satisfaction "
               "et décision doit être relu et validé par le Responsable Qualité (" +
               validator_role + ") avant présentation ou diffusion.")
    r1.font.size = Pt(15); r1.font.color.rgb = DARK
    p2 = tf.add_paragraph(); r2 = p2.add_run()
    r2.text = ("DRAFT, not a validated quality record. Every figure, finding, and "
               "decision must be validated by the RQ before use.")
    r2.font.size = Pt(12); r2.font.italic = True; r2.font.color.rgb = MUTED
    _add_footer(slide, "GEPROMED — Revue de direction (projet) · à valider par le RQ", logo_path)
    return slide


def build(content: dict, out_path: Path, logo_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    _title_slide(prs, content, logo_path)
    # Agenda slide from section headings.
    sections = content.get("sections") or [{"heading": h, "bullets": b} for h, b in DEFAULT_SECTIONS]
    agenda = [s.get("heading", "[section]") for s in sections]
    _content_slide(prs, "Ordre du jour / Agenda", agenda, logo_path)
    # One slide per section.
    for s in sections:
        _content_slide(prs, s.get("heading", "[section]"), s.get("bullets", []), logo_path)
    # Decisions & actions output slide.
    _actions_slide(prs, content.get("decisions"), content.get("actions"), logo_path)
    # Closing / validation.
    _closing_slide(prs, content.get("validator_role", "RQ"), logo_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def main() -> int:
    p = argparse.ArgumentParser(description="Generate the GEPROMED ISO 9001 §9.3 review deck (.pptx).")
    p.add_argument("--content", help="Path to a JSON content file (see --print-schema).")
    p.add_argument("--demo", action="store_true", help="Build a bracketed demo deck (no real figures).")
    p.add_argument("--out", default="revue_direction.pptx", help="Output .pptx path.")
    p.add_argument("--logo", default=str(DEFAULT_LOGO), help="Path to the GEPROMED logo PNG.")
    p.add_argument("--print-schema", action="store_true", help="Print the expected JSON schema and exit.")
    args = p.parse_args()

    if args.print_schema:
        print(json.dumps(SCHEMA, ensure_ascii=False, indent=2))
        return 0

    if args.demo:
        content = demo_content()
    elif args.content:
        path = Path(args.content)
        if not path.exists():
            print(f"ERROR: content file not found: {path}", file=sys.stderr)
            return 2
        try:
            content = json.loads(path.read_text(encoding="utf-8"))
        except json.JSONDecodeError as e:
            print(f"ERROR: invalid JSON in {path}: {e}", file=sys.stderr)
            return 2
    else:
        print("ERROR: provide --content <file.json> or --demo (or --print-schema).", file=sys.stderr)
        return 2

    logo = Path(args.logo)
    out = build(content, Path(args.out), logo)
    n_slides = len(Presentation(str(out)).slides._sldIdLst)
    print(f"Deck written: {out}  ({n_slides} slides)")
    print("Reminder: figures are bracketed placeholders — the RQ validates before use.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
